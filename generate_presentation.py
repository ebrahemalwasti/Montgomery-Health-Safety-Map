"""
generate_presentation.py
Generates a professional 12-slide PowerPoint presentation for the
WWV 2026 Hackathon submission of Montgomery Health & Safety Map.
Run once to produce: Montgomery_Health_Safety_Map_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
DARK_NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE  = RGBColor(0x93, 0xC5, 0xFD)
SAFETY_GREEN = RGBColor(0x16, 0xA3, 0x4A)
WARNING_AMBER = RGBColor(0xD9, 0x77, 0x06)
DANGER_RED  = RGBColor(0xDC, 0x26, 0x26)
LIGHT_GREY  = RGBColor(0xF1, 0xF5, 0xF9)
MID_NAVY    = RGBColor(0x1E, 0x3A, 0x5F)

# Slide dimensions (widescreen 13.333 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor,
             line_color=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Calibri", wrap=True):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_title_bar(slide, title_text, bg_color=DARK_NAVY, text_color=WHITE,
                  font_size=28):
    """Add a full-width title bar at the top of a content slide."""
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), bg_color)
    tf = bar.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = title_text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    # Vertical centring
    tf.margin_top = Inches(0.15)


def add_bullet_body(slide, bullets, left, top, width, height,
                    font_size=15, color=DARK_NAVY, bullet_char="\u2022 "):
    """Add a bulleted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.space_before = Pt(4)
        run = para.add_run()
        run.text = bullet_char + bullet
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tf


def add_callout_box(slide, text, left, top, width, height,
                    bg=ACCENT_BLUE, fg=WHITE, font_size=16, bold=True):
    """Add a highlighted callout rectangle with centred text."""
    shape = add_rect(slide, left, top, width, height, bg)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = fg
    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_NAVY)

    # Decorative accent bar
    add_rect(slide, Inches(0), Inches(2.9), SLIDE_W, Inches(0.06), ACCENT_BLUE)

    # Main title
    add_textbox(slide, "Montgomery Health & Safety Map",
                Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.5),
                font_size=44, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, "AI-Powered Civic Intelligence for a Safer City",
                Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.7),
                font_size=22, bold=False, color=LIGHT_BLUE,
                align=PP_ALIGN.CENTER)

    # Hackathon line
    add_textbox(slide, "WWV 2026 Hackathon  |  Team: Ebrahem Alwasti",
                Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
                font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Decorative bottom bar
    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), MID_NAVY)

    # Live demo URL bottom-right
    add_textbox(slide,
                "Live Demo: ebrahemalwasti.github.io/Montgomery-Health-Safety-Map/",
                Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45),
                font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)


def slide2_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "The Problem: Invisible Urban Safety Risks")

    bullets = [
        "Montgomery generates 25,000+ emergency calls per month across fragmented systems",
        "13+ civic datasets exist in isolation — 911 calls, traffic reports, health inspections, population data",
        "No tool connects infrastructure failures with healthcare access gaps at neighborhood level",
        "City planners lack block-by-block risk visibility for resource allocation",
        "Residents have zero awareness of their neighborhood's safety profile",
    ]
    add_bullet_body(slide, bullets,
                    Inches(0.5), Inches(1.25), Inches(12.3), Inches(4.5),
                    font_size=16)

    add_callout_box(slide, "The data exists. The insight doesn't.",
                    Inches(2.5), Inches(6.0), Inches(8.3), Inches(0.75),
                    font_size=18)


def slide3_objective(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "Our Mission: Predictive Civic Intelligence")

    bullets = [
        "Score every 500m \u00d7 500m zone in Montgomery on a 0\u2013100 Health & Safety scale",
        "Predict future risk trajectories using trend analysis \u2014 flag deteriorating zones before crises",
        "Identify resource gaps: pharmacy deserts, shelter blind spots, infrastructure decay",
        "Generate AI-written intelligence narratives for every neighborhood",
        "Pathfind citizens to nearest emergency resources in real time",
    ]
    add_bullet_body(slide, bullets,
                    Inches(0.5), Inches(1.25), Inches(12.3), Inches(4.5),
                    font_size=16)

    add_textbox(slide,
                "Aligned with WWV 2026 Challenge Areas: Smart City Tech  |  Public Health  |  Civic Innovation",
                Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
                font_size=13, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
                bold=True)


def slide4_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "Data: 13 Datasets from Montgomery Open Data Portal",
                  font_size=26)

    rows = [
        ("Dataset", "Used For"),
        ("911 Calls (175K+ records)", "Emergency pressure scoring"),
        ("Traffic Engineering Requests", "Infrastructure risk assessment"),
        ("Daily Population Trends", "Population vulnerability factor"),
        ("Food Safety Scores", "Food & environmental safety"),
        ("Fire Stations", "Proximity analysis"),
        ("Police Facilities", "Proximity analysis"),
        ("Tornado Shelters", "Emergency preparedness"),
        ("Pharmacy Locator", "Healthcare access / desert detection"),
        ("Community Centers", "Resource access scoring"),
        ("Parks & Trails", "Green space access"),
        ("Points of Interest", "Supplementary context"),
        ("City Limit Boundary", "Map boundary overlay"),
        ("Nuisance Complaints", "Environmental quality factor"),
    ]

    col1_w = Inches(5.5)
    col2_w = Inches(6.5)
    row_h = Inches(0.36)
    left = Inches(0.5)
    top_start = Inches(1.2)

    for i, (d, u) in enumerate(rows):
        top = top_start + i * row_h
        is_header = (i == 0)
        bg = DARK_NAVY if is_header else (MID_NAVY if i % 2 == 0 else ACCENT_BLUE)
        fg = WHITE

        add_rect(slide, left, top, col1_w, row_h, bg)
        add_rect(slide, left + col1_w, top, col2_w, row_h, bg)

        # Dataset cell
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.04),
                                      col1_w - Inches(0.2), row_h)
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = d
        r.font.name = "Calibri"
        r.font.size = Pt(11 if not is_header else 13)
        r.font.bold = is_header
        r.font.color.rgb = fg

        # Used For cell
        tb2 = slide.shapes.add_textbox(left + col1_w + Inches(0.1),
                                       top + Inches(0.04),
                                       col2_w - Inches(0.2), row_h)
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = u
        r2.font.name = "Calibri"
        r2.font.size = Pt(11 if not is_header else 13)
        r2.font.bold = is_header
        r2.font.color.rgb = fg

    add_textbox(slide,
                "Enhanced with Bright Data web enrichment \u2014 community sentiment from news & Yelp ratings",
                Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                font_size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
                bold=True)


def slide5_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "Methodology: Weighted Composite Scoring")

    # Formula box
    add_callout_box(slide,
                    "Score = 0.40 \u00d7 Emergency + 0.25 \u00d7 Infrastructure"
                    " + 0.20 \u00d7 Population + 0.15 \u00d7 Food Safety",
                    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75),
                    bg=DARK_NAVY, fg=WHITE, font_size=17, bold=True)

    components = [
        (DANGER_RED,     "Emergency Pressure (40%)",
         "911 call volumes, trend direction, seasonal patterns"),
        (WARNING_AMBER,  "Infrastructure Risk (25%)",
         "Broken signals, missing signs, dark intersections"),
        (ACCENT_BLUE,    "Population Vulnerability (20%)",
         "Daily flux patterns, density analysis"),
        (SAFETY_GREEN,   "Food & Environmental Safety (15%)",
         "Restaurant inspections, nuisance complaints"),
    ]

    box_w = Inches(2.9)
    box_h = Inches(1.5)
    gap = Inches(0.18)
    start_left = Inches(0.5)
    top = Inches(2.1)

    for idx, (color, label, desc) in enumerate(components):
        left = start_left + idx * (box_w + gap)
        shape = add_rect(slide, left, top, box_w, box_h, color)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label
        r1.font.name = "Calibri"
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = "Calibri"
        r2.font.size = Pt(11)
        r2.font.color.rgb = WHITE

    add_bullet_body(slide,
                    ["Spatial proximity via Haversine distance with bounding-box pre-filtering",
                     "Linear regression trend detection for Future Risk Score prediction",
                     "Automated AI narrative generation per zone"],
                    Inches(0.5), Inches(3.75), Inches(12.3), Inches(1.8),
                    font_size=14)


def slide6_map(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "The Map: A Geospatial Command Center")

    bullets = [
        "Score Map \u2014 Continuous color gradient: red (danger) \u2192 yellow \u2192 green (safe)",
        "Risk Map \u2014 Categorical view: High / Medium / Low risk classification",
        "Hotspot Detection \u2014 Highlights bottom 15% most dangerous zones",
        "Click-for-Details \u2014 Full AI analysis popup with risk drivers per zone",
        "City Limit Overlay \u2014 Municipal boundary for geographic context",
        "Zoom to High Risk \u2014 One-click auto-fit to worst neighborhoods",
        "Filter Mode \u2014 Isolate only high-risk areas for focused analysis",
    ]
    add_bullet_body(slide, bullets,
                    Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.0),
                    font_size=15)

    add_textbox(slide,
                "Built with MapLibre GL JS v4.1.3 \u2014 high-performance vector rendering, zero backend",
                Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                font_size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
                bold=True)


def slide7_predictive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "Predictive Intelligence: See Tomorrow\u2019s Risk Today")

    sections = [
        (DANGER_RED,    "Future Risk Scoring",
         ["Linear regression on monthly 911 call trends per zone",
          "Zones trending downward at >2 pts/month flagged for early intervention",
          "Example: Zone G4-22 scores 55 today \u2192 predicted 38 in 6 months"]),
        (ACCENT_BLUE,   "Spatial Proximity Analysis (Turf.js)",
         ["Nearest tornado shelter pathfinding with dashed-line visualization",
          "Nearest pharmacy distance calculation via Haversine formula",
          "Distance displayed in floating badge on map"]),
        (SAFETY_GREEN,  "Risk Classification",
         ["\U0001f534 High Risk: Score < 40",
          "\U0001f7e1 Medium Risk: Score 40\u201369",
          "\U0001f7e2 Low Risk: Score \u2265 70"]),
    ]

    box_w = Inches(3.9)
    box_h = Inches(3.5)
    gap = Inches(0.27)
    top = Inches(1.3)
    start_left = Inches(0.5)

    for idx, (color, heading, items) in enumerate(sections):
        left = start_left + idx * (box_w + gap)
        shape = add_rect(slide, left, top, box_w, box_h, color)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.1)
        # Heading
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = heading
        r0.font.name = "Calibri"
        r0.font.size = Pt(14)
        r0.font.bold = True
        r0.font.color.rgb = WHITE
        # Bullets
        for item in items:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            r = p.add_run()
            r.text = "\u2022 " + item
            r.font.name = "Calibri"
            r.font.size = Pt(12)
            r.font.color.rgb = WHITE


def slide8_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "Dashboard: Real-Time Civic Analytics")

    col_data = [
        (DARK_NAVY, "City Summary",
         ["Total zones scored",
          "City-wide average safety score",
          "High-risk zone count",
          "Low-risk zone count"]),
        (MID_NAVY, "Visualizations",
         ["SVG donut chart \u2014 risk distribution (High/Medium/Low)",
          "Bar charts \u2014 score distribution by range",
          "Critical zones list \u2014 click to fly to location",
          "Priority interventions ranked"]),
        (ACCENT_BLUE, "UX Features",
         ["Viewport-aware stats \u2014 analytics update as user pans",
          "Smooth hover tooltips with score + risk level",
          "Mobile-responsive layout \u2014 sidebar collapses on small screens",
          "Professional \u2018Civic Blue\u2019 design system"]),
    ]

    col_w = Inches(3.9)
    col_h = Inches(4.5)
    gap = Inches(0.27)
    top = Inches(1.3)

    for idx, (color, heading, items) in enumerate(col_data):
        left = Inches(0.5) + idx * (col_w + gap)
        shape = add_rect(slide, left, top, col_w, col_h, color)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.1)
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = heading
        r0.font.name = "Calibri"
        r0.font.size = Pt(15)
        r0.font.bold = True
        r0.font.color.rgb = WHITE
        for item in items:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            r = p.add_run()
            r.text = "\u2022 " + item
            r.font.name = "Calibri"
            r.font.size = Pt(12)
            r.font.color.rgb = WHITE


def slide9_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "Impact: Making Safety Visible and Equitable")

    stats = [
        (DANGER_RED,    "23 High-Risk Zones Identified"),
        (WARNING_AMBER, "14 Pharmacy Deserts Detected"),
        (WARNING_AMBER, "7 Shelter Coverage Gaps Mapped"),
        (DANGER_RED,    "34% of Zones Trending Worse"),
    ]

    box_w = Inches(2.9)
    box_h = Inches(0.75)
    gap = Inches(0.2)
    top = Inches(1.25)

    for idx, (color, text) in enumerate(stats):
        left = Inches(0.5) + idx * (box_w + gap)
        add_callout_box(slide, text, left, top, box_w, box_h,
                        bg=color, fg=WHITE, font_size=13, bold=True)

    bullets = [
        "City planners can generate prioritized intervention lists in 30 seconds",
        "Residents can search their address and understand their safety profile",
        "Emergency managers can identify shelter coverage blind spots",
        "Equity lens: highest-risk zones overlap with historically underserved neighborhoods",
    ]
    add_bullet_body(slide, bullets,
                    Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.5),
                    font_size=15)

    add_callout_box(slide,
                    "This tool doesn\u2019t just visualize data \u2014 it enables decisions that save lives.",
                    Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.75),
                    bg=DARK_NAVY, fg=WHITE, font_size=16, bold=True)


def slide10_innovation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "What Makes This Different")

    items = [
        ("AI-Generated Narratives",
         "Every zone gets a human-readable explanation of its risk drivers"),
        ("Multi-Dataset Fusion",
         "13 sources synthesized into a single intelligence layer"),
        ("Predictive, Not Reactive",
         "Future Risk Score flags deteriorating zones before crises"),
        ("Spatial Pathfinding",
         "Real-time nearest-resource routing with Turf.js"),
        ("Bright Data Enhancement",
         "Web-scraped sentiment adds a dynamic data dimension"),
        ("Viewport Intelligence",
         "Stats update based on what\u2019s visible, not the whole city"),
        ("Zero Infrastructure",
         "Runs entirely in the browser, no backend costs"),
    ]

    top = Inches(1.25)
    item_h = Inches(0.68)

    for i, (bold_part, rest) in enumerate(items):
        t = top + i * item_h
        # Accent stripe
        add_rect(slide, Inches(0.4), t + Inches(0.08),
                 Inches(0.07), Inches(0.45), ACCENT_BLUE)

        txBox = slide.shapes.add_textbox(Inches(0.6), t,
                                         Inches(12.0), item_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        r1 = para.add_run()
        r1.text = bold_part + " \u2014 "
        r1.font.name = "Calibri"
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = DARK_NAVY
        r2 = para.add_run()
        r2.text = rest
        r2.font.name = "Calibri"
        r2.font.size = Pt(15)
        r2.font.bold = False
        r2.font.color.rgb = DARK_NAVY


def slide11_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GREY)

    add_title_bar(slide, "From Hackathon to Product")

    # Timeline
    add_textbox(slide, "Roadmap", Inches(0.5), Inches(1.2), Inches(6.0),
                Inches(0.4), font_size=16, bold=True, color=DARK_NAVY)

    phases = [
        (SAFETY_GREEN, "Phase 1 (Now)",
         "Static pipeline + browser dashboard \u2014 COMPLETE"),
        (ACCENT_BLUE,  "Phase 2 (3 months)",
         "Automated API data ingestion + vector tile server"),
        (WARNING_AMBER, "Phase 3 (6 months)",
         "ML predictions (XGBoost) replacing heuristic model"),
        (DANGER_RED,   "Phase 4 (12 months)",
         "Multi-city SaaS platform for city planners"),
    ]

    for idx, (color, phase, desc) in enumerate(phases):
        top = Inches(1.65) + idx * Inches(0.88)
        add_rect(slide, Inches(0.5), top, Inches(0.15), Inches(0.6), color)
        txBox = slide.shapes.add_textbox(Inches(0.8), top,
                                         Inches(5.5), Inches(0.75))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = phase + ": "
        r1.font.name = "Calibri"
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = color
        r2 = p.add_run()
        r2.text = desc
        r2.font.name = "Calibri"
        r2.font.size = Pt(14)
        r2.font.color.rgb = DARK_NAVY

    # Revenue model
    add_rect(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(4.8), DARK_NAVY)
    txBox = slide.shapes.add_textbox(Inches(7.15), Inches(1.3),
                                     Inches(5.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = "Revenue Model"
    r0.font.name = "Calibri"
    r0.font.size = Pt(16)
    r0.font.bold = True
    r0.font.color.rgb = WHITE

    rev_items = [
        "Municipal licensing: $5K\u2013$25K/year per city",
        "Consulting: Custom dashboards for MPOs and emergency agencies",
        "API access for civic tech integrators",
    ]
    for item in rev_items:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = "\u2022 " + item
        r.font.name = "Calibri"
        r.font.size = Pt(13)
        r.font.color.rgb = WHITE

    add_textbox(slide,
                "Scalable to any city with an open data portal",
                Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                font_size=13, color=ACCENT_BLUE,
                align=PP_ALIGN.CENTER, bold=True)


def slide12_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_NAVY)

    # Accent line
    add_rect(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(0.06), ACCENT_BLUE)
    add_rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.06), ACCENT_BLUE)

    # Main title
    add_textbox(slide,
                "Montgomery\u2019s Data Was Always There.\nWe Made It Speak.",
                Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.3),
                font_size=28, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Takeaways
    takeaways = [
        "13 datasets \u2192 1 unified intelligence platform",
        "Every neighborhood scored, explained, and predicted",
        "Pharmacy deserts, shelter gaps, and infrastructure decay \u2014 all visible",
        "Zero backend. Zero cost. Maximum civic impact.",
    ]
    add_bullet_body(slide, takeaways,
                    Inches(1.5), Inches(1.7), Inches(10.3), Inches(2.0),
                    font_size=15, color=LIGHT_BLUE)

    # Live Demo label
    add_textbox(slide, "Live Demo",
                Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.55),
                font_size=26, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # URL
    add_textbox(slide,
                "https://ebrahemalwasti.github.io/Montgomery-Health-Safety-Map/",
                Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                font_size=18, color=LIGHT_BLUE,
                align=PP_ALIGN.CENTER)

    # GitHub
    add_textbox(slide,
                "github.com/ebrahemalwasti/Montgomery-Health-Safety-Map",
                Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.45),
                font_size=15, color=LIGHT_BLUE,
                align=PP_ALIGN.CENTER)

    # Thank you
    add_textbox(slide, "Thank you.  Questions?",
                Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
                font_size=22, bold=True, color=ACCENT_BLUE,
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide1_title(prs)
    slide2_problem(prs)
    slide3_objective(prs)
    slide4_data(prs)
    slide5_methodology(prs)
    slide6_map(prs)
    slide7_predictive(prs)
    slide8_dashboard(prs)
    slide9_impact(prs)
    slide10_innovation(prs)
    slide11_roadmap(prs)
    slide12_closing(prs)

    out = "Montgomery_Health_Safety_Map_Presentation.pptx"
    prs.save(out)
    print(f"\u2705 Presentation saved: {out}")


if __name__ == "__main__":
    main()
